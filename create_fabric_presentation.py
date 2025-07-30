#!/usr/bin/env python3
"""
Microsoft Fabric PowerPoint Presentation Generator
Creates a comprehensive presentation about Microsoft Fabric
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_fabric_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Define Microsoft brand colors
    ms_blue = RGBColor(0, 120, 212)
    ms_dark_blue = RGBColor(0, 78, 146)
    ms_light_blue = RGBColor(204, 229, 255)
    ms_gray = RGBColor(96, 94, 92)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Microsoft Fabric"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Unified Analytics Platform for the Modern Enterprise\n\nA Complete Data and Analytics Solution"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = ms_gray
    
    # Slide 2: What is Microsoft Fabric?
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What is Microsoft Fabric?"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Microsoft Fabric is an all-in-one analytics solution that covers everything from data movement to data science, Real-Time Analytics, and business intelligence.

Key Characteristics:
• Unified platform bringing together multiple analytics experiences
• Software as a Service (SaaS) offering
• Built on a foundation of compute and storage separated by design
• Integrates Power BI, Azure Synapse, and Azure Data Factory capabilities
• Provides a single, integrated environment for data professionals"""
    
    # Format content
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 3: Core Components
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Components of Microsoft Fabric"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Data Factory
• Data integration and ETL/ELT processes
• 200+ native connectors
• Visual data pipeline creation

Synapse Data Engineering
• Apache Spark-based big data processing
• Notebooks and job definitions
• Lakehouse architecture support

Synapse Data Warehouse
• SQL-based data warehousing
• Automatic optimization and scaling
• T-SQL compatibility

Synapse Data Science
• Machine learning model development
• MLflow integration
• Automated ML capabilities

Synapse Real-Time Analytics
• Real-time data ingestion and analysis
• KQL (Kusto Query Language) support
• Event streaming capabilities

Power BI
• Business intelligence and reporting
• Interactive dashboards and visualizations
• Self-service analytics"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 4: Key Benefits
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Benefits of Microsoft Fabric"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Unified Experience
• Single workspace for all analytics needs
• Consistent user interface across all workloads
• Simplified data governance and security

Simplified Architecture
• No need to piece together different services
• Built-in integration between components
• Reduced complexity and maintenance overhead

Cost Optimization
• Pay-as-you-go pricing model
• Automatic scaling and resource optimization
• Shared compute and storage resources

Enhanced Collaboration
• Shared workspace for data teams
• Git integration for version control
• Role-based access control

Future-Ready Platform
• Regular updates and new features
• AI and machine learning integration
• Support for emerging data formats and sources"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 5: OneLake - The Foundation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "OneLake: The Data Foundation"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """What is OneLake?
• Unified data lake built into Microsoft Fabric
• Single source of truth for all organizational data
• Automatically provisioned with every Fabric tenant

Key Features:
• Multi-format support (Delta Lake, Parquet, CSV, JSON)
• Hierarchical namespace organization
• Built-in data governance and lineage
• Automatic data discovery and cataloging
• Integration with Microsoft Purview

Benefits:
• Eliminates data silos
• Reduces data movement and duplication
• Provides consistent data access across all Fabric workloads
• Enables true self-service analytics
• Simplifies data management and governance"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 6: Use Cases and Scenarios
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Common Use Cases and Scenarios"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Enterprise Data Warehousing
• Modernize legacy data warehouse solutions
• Implement medallion architecture (Bronze, Silver, Gold)
• Support both batch and real-time data processing

Real-Time Analytics
• Monitor IoT devices and sensors
• Fraud detection and prevention
• Customer behavior analysis in real-time

Data Science and ML
• Build and deploy machine learning models
• Predictive analytics and forecasting
• Automated model training and deployment

Business Intelligence
• Self-service analytics for business users
• Executive dashboards and KPI monitoring
• Departmental reporting and analysis

Data Integration
• Migrate data from on-premises systems
• Integrate SaaS applications and cloud services
• Create unified views of customer and operational data"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 7: Getting Started
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Getting Started with Microsoft Fabric"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Prerequisites
• Microsoft 365 or Azure subscription
• Power BI Pro or Premium Per User license
• Fabric capacity (F64 or higher recommended for production)

First Steps:
1. Enable Fabric in your Power BI tenant settings
2. Create a new Fabric workspace
3. Choose your starting workload (Data Factory, Synapse, Power BI)
4. Begin with a pilot project or proof of concept

Best Practices:
• Start with a clear data strategy and governance plan
• Implement proper security and access controls
• Train your team on the new unified experience
• Leverage Microsoft's documentation and learning resources
• Consider engaging with Microsoft partners for implementation support

Resources:
• Microsoft Learn training modules
• Fabric documentation and samples
• Community forums and user groups
• Microsoft FastTrack for Fabric program"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 8: Pricing and Licensing
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Pricing and Licensing Model"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Capacity-Based Pricing
• Fabric uses a capacity-based pricing model
• Measured in Fabric Capacity Units (CUs)
• Pay for what you use with automatic scaling

Capacity Tiers:
• F2 (2 CUs) - Trial and development
• F4 (4 CUs) - Small workloads
• F8-F64 - Production workloads
• F128+ - Enterprise-scale deployments

What's Included:
• All Fabric workloads (Data Factory, Synapse, Power BI)
• OneLake storage (up to capacity limits)
• Compute resources for all analytics workloads
• Built-in security and governance features

Cost Optimization Tips:
• Use pause/resume for non-production workloads
• Implement data lifecycle management
• Monitor capacity utilization regularly
• Consider reserved capacity for predictable workloads"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 9: Roadmap and Future
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Microsoft Fabric Roadmap and Future"
    title.text_frame.paragraphs[0].font.color.rgb = ms_blue
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text = """Current Focus Areas:
• Enhanced AI and machine learning capabilities
• Improved performance and scalability
• Additional data connectors and integrations
• Advanced security and compliance features

Upcoming Features:
• Copilot integration across all workloads
• Enhanced real-time analytics capabilities
• Improved data visualization and reporting
• Better integration with Microsoft 365 ecosystem

Long-term Vision:
• Democratize data and analytics for all users
• Enable citizen data scientists and analysts
• Provide intelligent, automated insights
• Support for emerging data types and sources
• Seamless hybrid and multi-cloud scenarios

Stay Updated:
• Microsoft Fabric blog and announcements
• Monthly feature updates and releases
• Community feedback and feature requests
• Public preview programs for new capabilities"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = ms_gray
    
    # Slide 10: Thank You / Questions
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_shape.text_frame
    title_frame.text = "Thank You!"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.color.rgb = ms_blue
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Questions & Discussion\n\nMicrosoft Fabric: Unifying Your Analytics Journey"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = ms_gray
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    filename = "Microsoft_Fabric_Presentation.pptx"
    prs.save(filename)
    print(f"✅ PowerPoint presentation saved as: {filename}")
    print(f"📊 Total slides created: {len(prs.slides)}")
    
    return filename

if __name__ == "__main__":
    print("🚀 Creating Microsoft Fabric PowerPoint presentation...")
    filename = create_fabric_presentation()
    print(f"🎉 Presentation ready! Open '{filename}' in PowerPoint or LibreOffice Impress.")